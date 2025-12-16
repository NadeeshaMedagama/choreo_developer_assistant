"""
Text Extraction Services

Responsible for extracting text from various file formats.
Uses different extraction strategies based on file type.
"""

from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional
import io

from ..models import DiagramFile, ExtractedContent, FileType
from ..utils.logger import get_logger

logger = get_logger(__name__)


class TextExtractor(ABC):
    """Abstract base class for text extraction strategies."""

    @abstractmethod
    def can_extract(self, file: DiagramFile) -> bool:
        """Check if this extractor can handle the file type."""
        pass

    @abstractmethod
    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from the file."""
        pass


class ImageOCRExtractor(TextExtractor):
    """Extract text from images using OCR (Tesseract and Google Vision)."""

    def __init__(self, use_google_vision: bool = True, google_credentials: Optional[str] = None):
        """
        Initialize OCR extractor.

        Args:
            use_google_vision: Whether to use Google Vision API (more accurate)
            google_credentials: Path to Google Cloud credentials JSON
        """
        self.use_google_vision = use_google_vision
        self.google_credentials = google_credentials

        # Initialize Google Vision if available
        self.vision_client = None
        if use_google_vision:
            try:
                from google.cloud import vision
                import os
                if google_credentials:
                    os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = google_credentials
                self.vision_client = vision.ImageAnnotatorClient()
                logger.info("✓ Google Vision API initialized")
            except Exception as e:
                logger.warning(f"Could not initialize Google Vision API: {e}")
                logger.info("Falling back to Tesseract OCR")

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is an image type."""
        return file.file_type in [FileType.PNG, FileType.JPG, FileType.JPEG]

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from image using OCR."""
        logger.info(f"Extracting text from image: {file.file_name}")

        try:
            # Try Google Vision first
            if self.vision_client:
                return self._extract_with_google_vision(file)
            else:
                return self._extract_with_tesseract(file)
        except Exception as e:
            logger.error(f"OCR extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="ocr_failed",
                metadata={"error": str(e)}
            )

    def _extract_with_google_vision(self, file: DiagramFile) -> ExtractedContent:
        """Extract text using Google Vision API."""
        from google.cloud import vision

        with open(file.file_path, 'rb') as image_file:
            content = image_file.read()

        image = vision.Image(content=content)
        response = self.vision_client.document_text_detection(image=image)


        if response.error.message:
            raise Exception(response.error.message)

        # Get full text annotation
        text = response.full_text_annotation.text if response.full_text_annotation else ""

        # Calculate confidence (average of all word confidences)
        confidence = 0.0
        if response.text_annotations:
            confidences = [page.confidence for page in response.full_text_annotation.pages
                          if hasattr(page, 'confidence')]
            confidence = sum(confidences) / len(confidences) if confidences else 0.0

        logger.info(f"✓ Extracted {len(text)} characters (confidence: {confidence:.2%})")

        return ExtractedContent(
            source_file=file,
            raw_text=text,
            confidence_score=confidence,
            extraction_method="google_vision",
            metadata={
                "pages": len(response.full_text_annotation.pages) if response.full_text_annotation else 0
            }
        )

    def _extract_with_tesseract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text using Tesseract OCR."""
        try:
            import pytesseract
            from PIL import Image

            image = Image.open(file.file_path)

            # Get text and confidence
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
            text = pytesseract.image_to_string(image)

            # Calculate average confidence
            confidences = [int(conf) for conf in data['conf'] if conf != '-1']
            avg_confidence = sum(confidences) / len(confidences) / 100 if confidences else 0.0

            logger.info(f"✓ Extracted {len(text)} characters (confidence: {avg_confidence:.2%})")

            return ExtractedContent(
                source_file=file,
                raw_text=text,
                confidence_score=avg_confidence,
                extraction_method="tesseract",
                metadata={"words_detected": len([c for c in data['conf'] if c != '-1'])}
            )
        except Exception as e:
            logger.error(f"Tesseract extraction failed: {e}")
            raise


class SVGTextExtractor(TextExtractor):
    """Extract text from SVG files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is SVG."""
        return file.file_type == FileType.SVG

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from SVG file."""
        logger.info(f"Extracting text from SVG: {file.file_name}")

        try:
            from bs4 import BeautifulSoup

            with open(file.file_path, 'r', encoding='utf-8') as f:
                svg_content = f.read()

            soup = BeautifulSoup(svg_content, 'xml')

            # Extract all text elements
            text_elements = soup.find_all('text')
            texts = [elem.get_text(strip=True) for elem in text_elements if elem.get_text(strip=True)]

            # Also extract title and desc elements
            titles = [elem.get_text(strip=True) for elem in soup.find_all('title')]
            descs = [elem.get_text(strip=True) for elem in soup.find_all('desc')]

            all_text = '\n'.join(titles + descs + texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {len(texts)} text elements")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="svg_parsing",
                metadata={
                    "text_elements": len(texts),
                    "title_elements": len(titles),
                    "desc_elements": len(descs)
                }
            )
        except Exception as e:
            logger.error(f"SVG extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="svg_failed",
                metadata={"error": str(e)}
            )


class DrawioExtractor(TextExtractor):
    """Extract text from Draw.io files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is drawio."""
        return file.file_type == FileType.DRAWIO

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from Draw.io file."""
        logger.info(f"Extracting text from Draw.io: {file.file_name}")

        try:
            from bs4 import BeautifulSoup
            import urllib.parse

            with open(file.file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            soup = BeautifulSoup(content, 'xml')

            # Draw.io stores diagram data in mxCell elements
            cells = soup.find_all('mxCell')
            texts = []

            for cell in cells:
                # Get value attribute (contains text)
                value = cell.get('value', '')
                if value:
                    # Decode HTML entities
                    decoded = urllib.parse.unquote(value)
                    if decoded.strip():
                        texts.append(decoded.strip())

            all_text = '\n'.join(texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {len(texts)} cells")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="drawio_parsing",
                metadata={"cells_processed": len(texts)}
            )
        except Exception as e:
            logger.error(f"Draw.io extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="drawio_failed",
                metadata={"error": str(e)}
            )


class PDFExtractor(TextExtractor):
    """Extract text from PDF files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is PDF."""
        return file.file_type == FileType.PDF

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from PDF."""
        logger.info(f"Extracting text from PDF: {file.file_name}")

        try:
            import PyPDF2

            texts = []
            with open(file.file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = len(pdf_reader.pages)

                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        texts.append(text)

            all_text = '\n\n'.join(texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {num_pages} pages")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="pdf_parsing",
                metadata={"pages": num_pages}
            )
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="pdf_failed",
                metadata={"error": str(e)}
            )


class DocxExtractor(TextExtractor):
    """Extract text from DOCX files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is DOCX."""
        return file.file_type == FileType.DOCX

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from DOCX."""
        logger.info(f"Extracting text from DOCX: {file.file_name}")

        try:
            from docx import Document

            doc = Document(file.file_path)
            texts = [para.text for para in doc.paragraphs if para.text.strip()]
            all_text = '\n\n'.join(texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {len(texts)} paragraphs")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="docx_parsing",
                metadata={"paragraphs": len(texts)}
            )
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="docx_failed",
                metadata={"error": str(e)}
            )


class XlsxExtractor(TextExtractor):
    """Extract text from XLSX files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is XLSX."""
        return file.file_type == FileType.XLSX

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from XLSX."""
        logger.info(f"Extracting text from XLSX: {file.file_name}")

        try:
            import openpyxl

            workbook = openpyxl.load_workbook(file.file_path, data_only=True)
            texts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                texts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip(' |'):
                        texts.append(row_text)

            all_text = '\n'.join(texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {len(workbook.sheetnames)} sheets")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="xlsx_parsing",
                metadata={"sheets": len(workbook.sheetnames)}
            )
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="xlsx_failed",
                metadata={"error": str(e)}
            )


class PptxExtractor(TextExtractor):
    """Extract text from PPTX (PowerPoint) files."""

    def can_extract(self, file: DiagramFile) -> bool:
        """Check if file is PPTX."""
        return file.file_type == FileType.PPTX

    def extract(self, file: DiagramFile) -> ExtractedContent:
        """Extract text from PPTX."""
        logger.info(f"Extracting text from PPTX: {file.file_name}")

        try:
            from pptx import Presentation

            prs = Presentation(file.file_path)
            texts = []
            slide_count = 0

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                slide_texts = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                    # Also check for text in tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    slide_texts.append(cell.text.strip())

                if slide_texts:
                    texts.append(f"=== Slide {slide_num} ===")
                    texts.extend(slide_texts)
                    texts.append("")  # Empty line between slides

            all_text = '\n'.join(texts)

            logger.info(f"✓ Extracted {len(all_text)} characters from {slide_count} slides")

            return ExtractedContent(
                source_file=file,
                raw_text=all_text,
                confidence_score=1.0,
                extraction_method="pptx_parsing",
                metadata={"slides": slide_count}
            )
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            return ExtractedContent(
                source_file=file,
                raw_text="",
                extraction_method="pptx_failed",
                metadata={"error": str(e)}
            )


class TextExtractionService:
    """
    Main service for text extraction.
    Uses appropriate extractor based on file type.
    """

    def __init__(self, google_credentials: Optional[str] = None):
        """
        Initialize text extraction service.

        Args:
            google_credentials: Path to Google Cloud credentials
        """
        self.extractors = [
            ImageOCRExtractor(use_google_vision=True, google_credentials=google_credentials),
            SVGTextExtractor(),
            DrawioExtractor(),
            PDFExtractor(),
            DocxExtractor(),
            XlsxExtractor(),
            PptxExtractor(),  # Add PPTX extractor
        ]

    def extract_text(self, file: DiagramFile) -> ExtractedContent:
        """
        Extract text from a file using appropriate extractor.

        Args:
            file: DiagramFile to process

        Returns:
            ExtractedContent with text and metadata
        """
        for extractor in self.extractors:
            if extractor.can_extract(file):
                return extractor.extract(file)

        logger.warning(f"No extractor found for file type: {file.file_type}")
        return ExtractedContent(
            source_file=file,
            raw_text="",
            extraction_method="unsupported",
            metadata={"error": f"Unsupported file type: {file.file_type}"}
        )
